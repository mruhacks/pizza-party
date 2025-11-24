"""
PowerPoint to Markdown Converter
Extracts slide content and presenter notes from PPTX files and outputs formatted Markdown.
"""

import sys
from pathlib import Path
from pptx import Presentation


def extract_slide_text(slide):
    """Extract all text content from a slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())
    return "\n\n".join(text_parts)


def extract_notes(slide):
    """Extract presenter notes from a slide."""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        return notes_text
    return ""


def pptx_to_markdown(pptx_path):
    """Convert a PowerPoint file to Markdown format."""
    prs = Presentation(pptx_path)
    markdown_lines = []
    
    # Add title
    pptx_name = Path(pptx_path).stem
    markdown_lines.append(f"# {pptx_name}\n")
    
    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        # Slide heading
        markdown_lines.append(f"## Slide {i}\n")
        
        # Slide content
        slide_text = extract_slide_text(slide)
        if slide_text:
            markdown_lines.append(slide_text + "\n")
        
        # Presenter notes
        notes = extract_notes(slide)
        if notes:
            markdown_lines.append("### Presenter Notes\n")
            markdown_lines.append(notes + "\n")
        
        # Add separator between slides
        markdown_lines.append("---\n")
    
    return "\n".join(markdown_lines)


def main():
    if len(sys.argv) < 2:
        print("Usage: python main.py <path-to-pptx-file> [output-file]")
        print("\nExample:")
        print("  python main.py presentation.pptx")
        print("  python main.py presentation.pptx output.md")
        sys.exit(1)
    
    input_path = sys.argv[1]
    
    # Validate input file
    if not Path(input_path).exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    if not input_path.lower().endswith('.pptx'):
        print("Error: Input file must be a .pptx file")
        sys.exit(1)
    
    # Determine output path
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        output_path = Path(input_path).stem + ".md"
    
    # Convert and save
    print(f"Converting {input_path} to Markdown...")
    markdown_content = pptx_to_markdown(input_path)
    
    Path(output_path).write_text(markdown_content, encoding='utf-8')
    print(f"✓ Saved to {output_path}")


if __name__ == "__main__":
    main()
